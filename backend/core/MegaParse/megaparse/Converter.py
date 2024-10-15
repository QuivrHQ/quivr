import asyncio
import logging
import os
from collections import Counter
from pathlib import Path
from typing import List, Set

import pandas as pd
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.text.run import Run
from langchain_community.document_loaders.base import BaseLoader
from langchain_core.documents import Document as LangChainDocument
from llama_index.core.schema import Document as LlamaDocument
from llama_parse import LlamaParse
from llama_parse.utils import Language, ResultType
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from megaparse.config import MegaparseConfig, PdfParser
from megaparse.markdown_processor import MarkdownProcessor
from megaparse.multimodal_convertor.megaparse_vision import MegaParseVision
from megaparse.unstructured_convertor import ModelEnum, UnstructuredParser

logger = logging.getLogger("megaparse")


class Converter:
    def __init__(self) -> None:
        pass

    async def convert(self, file_path: str | Path) -> LangChainDocument:
        raise NotImplementedError("Subclasses should implement this method")

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        with open(file_path, "w") as f:
            f.write(md_content)


class XLSXConverter(Converter):
    def __init__(self) -> None:
        pass

    async def convert(self, file_path: str | Path) -> LangChainDocument:
        if isinstance(file_path, str):
            file_path = Path(file_path)
        xls = pd.ExcelFile(file_path)  # type: ignore
        sheets = pd.read_excel(xls)

        target_text = self.table_to_text(sheets)

        return LangChainDocument(
            page_content=target_text,
            metadata={"filename": file_path.name, "type": "xlsx"},
        )

    def convert_tab(self, file_path: str | Path, tab_name: str) -> str:
        if isinstance(file_path, str):
            file_path = Path(file_path)
        xls = pd.ExcelFile(str(file_path))
        sheets = pd.read_excel(xls, tab_name)
        target_text = self.table_to_text(sheets)
        return target_text

    def table_to_text(self, df):
        text_rows = []
        for _, row in df.iterrows():
            row_text = " | ".join(str(value) for value in row.values if pd.notna(value))
            if row_text:
                text_rows.append("|" + row_text + "|")
        return "\n".join(text_rows)


class DOCXConverter(Converter):
    def __init__(self) -> None:
        self.header_handled = False

    async def convert(self, file_path: str | Path) -> LangChainDocument:
        if isinstance(file_path, str):
            file_path = Path(file_path)
        doc = Document(str(file_path))
        md_content = []
        # Handle header
        if doc.sections and doc.sections[0].header:
            header_content = self._handle_header(doc.sections[0].header)
            if header_content:
                md_content.append(header_content)

        for element in doc.element.body:
            if isinstance(element, CT_P):
                md_content.append(self._handle_paragraph(Paragraph(element, doc)))
            elif isinstance(element, CT_Tbl):
                md_content += self._handle_table(Table(element, doc))
            # Add more handlers here (image, header, footer, etc)

        return LangChainDocument(
            page_content="\n".join(md_content),
            metadata={"filename": file_path.name, "type": "docx"},
        )

    def _handle_header(self, header) -> str:
        if not self.header_handled:
            parts = []
            for paragraph in header.paragraphs:
                parts.append(f"# {paragraph.text}")
            for table in header.tables:
                parts += self._handle_header_table(table)
            self.header_handled = True
            return "\n".join(parts)
        return ""

    def _handle_header_table(self, table: Table) -> List[str]:
        cell_texts = [cell.text for row in table.rows for cell in row.cells]
        cell_texts.remove("")
        # Find the most repeated cell text
        text_counts = Counter(cell_texts)
        title = text_counts.most_common(1)[0][0] if cell_texts else ""
        other_texts = [text for text in cell_texts if text != title and text != ""]

        md_table_content = []
        if title:
            md_table_content.append(f"# {title}")
        for text in other_texts:
            md_table_content.append(f"*{text}*;")
        return md_table_content

    def _handle_paragraph(self, paragraph: Paragraph) -> str:
        if paragraph.style.name.startswith("Heading"):  # type: ignore
            level = int(paragraph.style.name.split()[-1])  # type: ignore
            return f"{'#' * level} {paragraph.text}"
        else:
            parts = []
            for run in paragraph.runs:
                if run.text != "":
                    parts.append(self._handle_run(run))
            return "".join(parts)

    def _handle_run(self, run: Run) -> str:
        text: str = run.text
        if run.bold:
            if len(text) < 200:
                # FIXME : handle table needs to be improved -> have the paragraph they are in
                text = f"## {text}"
            else:
                text = f"**{text}**"
        if run.italic:
            text = f"*{text}*"
        return text

    def _handle_table(self, table: Table) -> List[str]:
        row_content = []
        for i, row in enumerate(table.rows):
            row_content.append(
                "| " + " | ".join(cell.text.strip() for cell in row.cells) + " |"
            )
            if i == 0:
                row_content.append("|" + "---|" * len(row.cells))

        return row_content

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        with open(file_path, "w") as f:
            f.write(md_content)


class PPTXConverter:
    def __init__(self, add_images=False) -> None:
        self.header_handled = False
        self.add_images = add_images

    async def convert(self, file_path: str | Path) -> LangChainDocument:
        if isinstance(file_path, str):
            file_path = Path(file_path)
        prs = Presentation(str(file_path))
        md_content = []
        unique_slides: Set[str] = set()

        # Handle header
        if prs.slides and prs.slides[0].placeholders:
            header_content = self._handle_header(prs.slides[0].placeholders)
            if header_content:
                md_content.append(header_content)

        for i, slide in enumerate(prs.slides):
            slide_md_content: List[str] = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # type: ignore
                    slide_md_content += self._handle_table(shape.table)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.add_images:  # type: ignore
                    slide_md_content.append(self._handle_image(shape))
                elif hasattr(shape, "text"):
                    slide_md_content.append(self._handle_paragraph(shape.text))

            slide_md_str = "\n".join(slide_md_content)
            if slide_md_str not in unique_slides:
                unique_slides.add(slide_md_str)
                slide_md_str = f"## Slide {i+1}\n{slide_md_str}"
                md_content.append(slide_md_str)

        return LangChainDocument(
            page_content="\n".join(md_content),
            metadata={"filename": file_path.name, "type": "pptx"},
        )

    def _handle_header(self, placeholders) -> str:
        if not self.header_handled:
            parts = []
            for placeholder in placeholders:
                if placeholder.placeholder_format.idx == 0:  # Title placeholder
                    parts.append(f"# {placeholder.text}")
                elif placeholder.placeholder_format.idx == 1:  # Subtitle placeholder
                    parts.append(f"## {placeholder.text}")
            self.header_handled = True
            return "\n".join(parts)
        return ""

    def _handle_paragraph(self, text: str) -> str:
        # Assuming text is a simple paragraph without complex formatting
        # if text contains letters return text
        if any(c.isalpha() for c in text):
            return text + "\n"
        return ""

    def _handle_image(self, shape) -> str:
        image = shape.image
        image_bytes = image.blob
        image_format = image.ext
        image_filename = f"images/image_{shape.shape_id}.{image_format}"
        with open(image_filename, "wb") as f:
            f.write(image_bytes)
        return f"![Image {shape.shape_id}](../{image_filename})"

    def _handle_table(self, table) -> List[str]:
        row_content = []
        for i, row in enumerate(table.rows):
            row_content.append(
                "| " + " | ".join(cell.text.strip() for cell in row.cells) + " |"
            )
            if i == 0:
                row_content.append("|" + "---|" * len(row.cells))
        return row_content

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        with open(file_path, "w") as f:
            f.write(md_content)


class PDFConverter:
    def __init__(
        self,
        llama_parse_api_key: str,
        method: PdfParser | str = PdfParser.UNSTRUCTURED,
        model=ModelEnum.NONE,
        strategy="fast",
    ) -> None:
        self.strategy = strategy
        self.llama_parse_api_key = llama_parse_api_key
        if isinstance(method, str):
            try:
                method = PdfParser(method)
            except ValueError:
                raise ValueError(f"Method {method} not supported")
        self.method = method

    async def _llama_parse(self, api_key: str, file_path: str | Path):
        logger.debug(f"Parsing {file_path.name} using llama_parse")
        parsing_instructions = "Do not take into account the page breaks (no --- between pages), do not repeat the header and the footer so the tables are merged. Keep the same format for similar tables."
        self.parser = LlamaParse(
            api_key=str(api_key),
            result_type=ResultType.MD,
            gpt4o_mode=True,
            verbose=True,
            language=Language.FRENCH,
            parsing_instruction=parsing_instructions,  # Optionally you can define a parsing instruction
        )
        documents: List[LlamaDocument] = await self.parser.aload_data(str(file_path))
        parsed_md = ""
        for document in documents:
            text_content = document.text
            parsed_md = parsed_md + text_content
        return parsed_md

    def _unstructured_parse(
        self, file_path: str | Path, model: ModelEnum = ModelEnum.NONE
    ):
        logger.debug(
            f"Parsing {file_path.name} using unstructured with strategy {self.strategy}"
        )
        unstructured_parser = UnstructuredParser()
        return unstructured_parser.convert(
            file_path, model=model, strategy=self.strategy
        )

    async def _lmm_parse(self, file_path: str | Path):
        lmm_parser = MegaParseVision()
        return await lmm_parser.parse(file_path)

    async def convert(
        self,
        file_path: str | Path,
        model: ModelEnum = ModelEnum.NONE,
        gpt4o_cleaner=False,
    ) -> LangChainDocument:
        if isinstance(file_path, str):
            file_path = Path(file_path)

        parsed_md = ""
        if self.method == PdfParser.LLAMA_PARSE:
            assert (
                self.llama_parse_api_key is not None
            ), "LLama Parse API key is required for this method"
            parsed_md = await self._llama_parse(self.llama_parse_api_key, file_path)
        elif self.method == PdfParser.MEGAPARSE_VISION:
            parsed_md = await self._lmm_parse(file_path)
        elif self.method == PdfParser.UNSTRUCTURED:
            parsed_md = self._unstructured_parse(file_path, model)
        else:
            raise ValueError(f"Method {self.method} not supported")

        if gpt4o_cleaner:
            md_processor = MarkdownProcessor(
                parsed_md,
                strict=True,
                remove_pagination=True,
            )
            md_cleaned = md_processor.process(gpt4o_cleaner=gpt4o_cleaner)
            parsed_md = md_cleaned

        if (
            len(parsed_md) < 5
            and file_path.stat().st_size > 100
            and self.strategy == "fast"
        ):
            if os.environ.get("LLAMA_PARSE_API_KEY"):
                logger.info(f"Switching to llama parse strategy for {file_path.name}")
                self.method = PdfParser.LLAMA_PARSE
                self.llama_parse_api_key = os.environ.get("LLAMA_PARSE_API_KEY")
                return await self.convert(file_path, model, gpt4o_cleaner=gpt4o_cleaner)
            else:
                logger.info(
                    f"Unable to switch to llama parse strategy for {file_path.name}"
                )

        return LangChainDocument(
            page_content=parsed_md,
            metadata={"filename": file_path.name, "type": "pdf"},
        )

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        with open(file_path, "w") as f:
            f.write(md_content)


class MegaParse(BaseLoader):
    def __init__(
        self,
        file_path: str | Path,
        config: MegaparseConfig = MegaparseConfig(),
    ) -> None:
        if isinstance(file_path, str):
            file_path = Path(file_path)
        self.file_path = file_path
        self.config = config

    async def aload(self, **convert_kwargs) -> LangChainDocument:
        file_extension: str = os.path.splitext(self.file_path)[1]
        if file_extension == ".docx":
            converter = DOCXConverter()
        elif file_extension == ".pptx":
            converter = PPTXConverter()
        elif file_extension == ".pdf":
            converter = PDFConverter(
                llama_parse_api_key=str(self.config.llama_parse_api_key),
                strategy=self.config.strategy,
                method=self.config.pdf_parser,
            )
        elif file_extension == ".xlsx":
            converter = XLSXConverter()
        else:
            raise ValueError(f"Unsupported file extension: {file_extension}")

        return await converter.convert(self.file_path, **convert_kwargs)

    def load(self, **kwargs) -> LangChainDocument:
        file_extension: str = os.path.splitext(self.file_path)[1]
        if file_extension == ".docx":
            converter = DOCXConverter()
        elif file_extension == ".pptx":
            converter = PPTXConverter()
        elif file_extension == ".pdf":
            converter = PDFConverter(
                llama_parse_api_key=str(self.config.llama_parse_api_key),
                strategy=self.config.strategy,
            )
        elif file_extension == ".xlsx":
            converter = XLSXConverter()
        else:
            print(self.file_path, file_extension)
            raise ValueError(f"Unsupported file extension: {file_extension}")

        loop = asyncio.get_event_loop()
        return loop.run_until_complete(converter.convert(self.file_path, **kwargs))

    def load_tab(self, tab_name: str, **kwargs) -> LangChainDocument:
        file_extension: str = os.path.splitext(self.file_path)[1]
        if file_extension == ".xlsx":
            converter = XLSXConverter()
        else:
            print(self.file_path, file_extension)
            raise ValueError(f"Unsupported file extension for tabs: {file_extension}")

        result = converter.convert_tab(self.file_path, tab_name=tab_name)
        return LangChainDocument(
            page_content=result,
            metadata={"filename": self.file_path.name, "type": "xlsx"},
        )

    def save_md(self, md_content: str, file_path: Path | str) -> None:
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        with open(file_path, "w+") as f:
            f.write(md_content)
